"""Generate a PowerPoint deck summarizing what's new in Microsoft AI (February 2026)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colour palette (Microsoft-inspired) ──────────────────────────────────────
DARK_BG    = RGBColor(0x1B, 0x1B, 0x1B)
ACCENT     = RGBColor(0x00, 0x78, 0xD4)  # Microsoft blue
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ORANGE     = RGBColor(0xFF, 0x8C, 0x00)
GREEN      = RGBColor(0x10, 0x7C, 0x10)
PURPLE     = RGBColor(0x88, 0x61, 0xC9)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


# ── Helpers ───────────────────────────────────────────────────────────────────

def _set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_text_box(slide, left, top, width, height, text, font_size=18,
                  bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_bullet_slide(slide, items, start_top, left=Inches(0.8), width=Inches(11.5),
                      font_size=16, color=WHITE):
    """Add a text box with bullet list items."""
    txBox = slide.shapes.add_textbox(left, start_top, width, Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (title, desc) in enumerate(items):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.space_before = Pt(8)
        p.space_after = Pt(4)

        run_title = p.add_run()
        run_title.text = title
        run_title.font.size = Pt(font_size)
        run_title.font.bold = True
        run_title.font.color.rgb = ACCENT
        run_title.font.name = "Segoe UI"

        if desc:
            run_desc = p.add_run()
            run_desc.text = f"\n{desc}"
            run_desc.font.size = Pt(font_size - 2)
            run_desc.font.bold = False
            run_desc.font.color.rgb = color
            run_desc.font.name = "Segoe UI"
    return txBox


def _accent_bar(slide, top):
    _add_rect(slide, Inches(0.5), top, Inches(1.2), Pt(4), ACCENT)


def _style_cell(cell, text, bold=False, bg=None, fg=WHITE, size=13):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = fg
    run.font.name = "Segoe UI"
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


# ── Slide 1: Title ──────────────────────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
_set_bg(slide, DARK_BG)
_add_rect(slide, 0, 0, W, Inches(0.15), ACCENT)

_add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1.5),
              "What's New in Microsoft AI", font_size=44, bold=True, color=WHITE)
_add_text_box(slide, Inches(0.8), Inches(3.0), Inches(11), Inches(1),
              "Microsoft Foundry, Azure OpenAI Service & Related Platforms",
              font_size=24, color=LIGHT_GRAY)
_add_text_box(slide, Inches(0.8), Inches(4.2), Inches(11), Inches(0.8),
              "February 2026", font_size=20, color=ACCENT)
_add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.6),
              "Sources: Azure Blog, Microsoft Learn, Microsoft Tech Community",
              font_size=14, color=LIGHT_GRAY)


# ── Slide 2: Agenda / TOC ───────────────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide, DARK_BG)
_add_rect(slide, 0, 0, W, Inches(0.15), ACCENT)
_add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
              "Agenda", font_size=36, bold=True, color=WHITE)
_accent_bar(slide, Inches(1.3))

topics = [
    ("1.", "New Models in Microsoft Foundry"),
    ("2.", "AI Infrastructure & Hardware"),
    ("3.", "Azure Data Services for AI"),
    ("4.", "Agentic AI & Frameworks"),
    ("5.", "Foundry Blog Highlights"),
    ("6.", "Model Deprecations & Retirements"),
    ("7.", "Key February 2026 Quick-Reference Table"),
]
y = Inches(1.8)
for num, topic in topics:
    _add_text_box(slide, Inches(1.0), y, Inches(0.5), Inches(0.5), num,
                  font_size=20, bold=True, color=ACCENT)
    _add_text_box(slide, Inches(1.6), y, Inches(10), Inches(0.5), topic,
                  font_size=20, color=WHITE)
    y += Inches(0.65)


# ── Slide 3: New Models – Claude in Foundry ─────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide, DARK_BG)
_add_rect(slide, 0, 0, W, Inches(0.15), ACCENT)
_add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
              "Anthropic Claude Models in Microsoft Foundry", font_size=34, bold=True, color=WHITE)
_accent_bar(slide, Inches(1.1))

items = [
    ("Claude Opus 4.6 (Feb 5, 2026)",
     "Most powerful Claude model for coding, agents & enterprise workflows. "
     "Delegates complex tasks end-to-end with autonomous production execution."),
    ("Claude Sonnet 4.6 (Feb 17, 2026)",
     "Frontier performance across coding, agents & professional work at scale. "
     "Complements Opus for teams needing a balance of performance and throughput."),
    ("Why It Matters",
     "\u2022 Broadens the model catalog in Microsoft Foundry beyond OpenAI\n"
     "\u2022 Enterprise-grade security, compliance & governance through Foundry Control Plane\n"
     "\u2022 Enables customers to choose the best model per workload"),
]
_add_bullet_slide(slide, items, Inches(1.5))


# ── Slide 4: AI Infrastructure & Hardware ────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide, DARK_BG)
_add_rect(slide, 0, 0, W, Inches(0.15), ACCENT)
_add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
              "AI Infrastructure & Hardware", font_size=34, bold=True, color=WHITE)
_accent_bar(slide, Inches(1.1))

items = [
    ("High-Temperature Superconductors (Feb 10, 2026)",
     "Microsoft is exploring superconductor technology to transform datacenter power "
     "infrastructure, addressing rising demands from AI-intensive workloads."),
]
_add_bullet_slide(slide, items, Inches(1.5))


# ── Slide 5: Azure Data Services ────────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide, DARK_BG)
_add_rect(slide, 0, 0, W, Inches(0.15), ACCENT)
_add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
              "Azure Data Services for AI", font_size=34, bold=True, color=WHITE)
_accent_bar(slide, Inches(1.1))

items = [
    ("PostgreSQL on Azure \u2013 Supercharged for AI (Feb 2, 2026)",
     "\u2022 GitHub Copilot integration for SQL in VS Code\n"
     "\u2022 Invoke LLMs via SQL through Microsoft Foundry integration\n"
     "\u2022 DiskANN vector indexing for high-performance similarity search\n"
     "\u2022 MCP server connecting PostgreSQL \u2192 Foundry Agent Framework\n"
     "\u2022 Zero-ETL mirroring to Microsoft Fabric\n"
     "\u2022 PostgreSQL 18 GA, new V6 compute SKUs, Elastic Clusters"),
    ("Azure HorizonDB (Private Preview)",
     "New PostgreSQL-compatible service built for AI-native workloads. "
     "Scale-out compute, sub-millisecond latency, built-in AI features."),
]
_add_bullet_slide(slide, items, Inches(1.5))


# ── Slide 6: Agentic AI ─────────────────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide, DARK_BG)
_add_rect(slide, 0, 0, W, Inches(0.15), ACCENT)
_add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
              "Agentic AI & Frameworks", font_size=34, bold=True, color=WHITE)
_accent_bar(slide, Inches(1.1))

items = [
    ("Agentic Cloud Operations (Feb 11, 2026)",
     "New paradigm where AI agents assist with cloud operations \u2013 monitoring, "
     "incident response & management tasks powered by the Foundry platform."),
    ("Pantone \u2013 Agentic AI Case Study (Feb 12, 2026)",
     "Pantone launched an AI-powered experience combining an AI-ready database with "
     "agentic workflows, iterating rapidly on real user feedback."),
    ("Agents League Community Challenge (Feb 16\u201327, 2026)",
     "Developer challenge: build agents, compete live, win prizes. "
     "Part of Microsoft\u2019s push to grow the agentic AI developer ecosystem."),
]
_add_bullet_slide(slide, items, Inches(1.5))


# ── Slide 7: Foundry Blog Highlights ────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide, DARK_BG)
_add_rect(slide, 0, 0, W, Inches(0.15), ACCENT)
_add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
              "Foundry Blog Highlights (Feb 2026)", font_size=34, bold=True, color=WHITE)
_accent_bar(slide, Inches(1.1))

items = [
    ("Microsoft Agent Framework \u2013 Release Candidate (Feb 19)",
     "Successor to Semantic Kernel & AutoGen. Stable API for v1.0 in .NET & Python. "
     "Graph-based workflows, multi-provider support (Foundry, OpenAI, Claude, Bedrock, Ollama), "
     "A2A / AG-UI / MCP interoperability."),
    ("What\u2019s New in Foundry | Dec 2025 & Jan 2026 (Feb 18)",
     "Comprehensive roundup: GPT-5.2 GA, Codex Max GA, Mistral Large 3, DeepSeek V3.2, "
     "Kimi-K2 Thinking, Cohere Rerank 4, GPT-image-1.5 GA, FLUX.2 [pro], Memory in Agent Service, "
     "A2A Tool, Computer Use, Foundry MCP Server, azure-ai-projects v2 beta."),
    ("DPO Fine-Tuning Using Foundry SDK (Feb 13)",
     "Guide to Direct Preference Optimization \u2013 learn from human preference pairs "
     "to align model behaviour without a separate reward model."),
    ("Beyond the Prompt \u2013 Why & How to Fine-tune (Feb 6)",
     "When prompt engineering & RAG aren\u2019t enough: fine-tuning for consistent, "
     "policy-compliant outputs at enterprise scale."),
]
_add_bullet_slide(slide, items, Inches(1.5))


# ── Slide 8: Deprecations – Imminent Retirements ────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide, DARK_BG)
_add_rect(slide, 0, 0, W, Inches(0.15), ACCENT)
_add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
              "Model Deprecations & Retirements", font_size=34, bold=True, color=WHITE)
_accent_bar(slide, Inches(1.1))

# Table – imminent retirements
ret_rows, ret_cols = 7, 4
tbl_shape = slide.shapes.add_table(ret_rows, ret_cols, Inches(0.5), Inches(1.5),
                                    Inches(12.3), Inches(3.8))
ret_table = tbl_shape.table
ret_table.columns[0].width = Inches(3.8)
ret_table.columns[1].width = Inches(2.5)
ret_table.columns[2].width = Inches(3.2)
ret_table.columns[3].width = Inches(2.8)

ret_headers = ["Model", "Retirement", "Auto-Upgrade", "Upgrade Path"]
ret_data = [
    ["gpt-5-chat (preview)", "Mar 1, 2026", "\u2013", "gpt-5.2-chat"],
    ["gpt-4o Standard (all GA)", "Mar 31, 2026", "Starts Mar 9", "gpt-5.1"],
    ["gpt-4o-mini Standard", "Mar 31, 2026", "Starts Mar 9", "gpt-4.1-mini"],
    ["o1-pro", "Sep 18, 2026", "\u2013", "o3-pro"],
    ["gpt-4.1 family", "Oct 14, 2026", "\u2013", "gpt-5 / gpt-5-mini / gpt-5-nano"],
    ["o3 / o4-mini", "Oct 16, 2026", "\u2013", "TBD"],
]

for i, h in enumerate(ret_headers):
    _style_cell(ret_table.cell(0, i), h, bold=True, bg=ACCENT, fg=WHITE, size=14)

row_bg = [RGBColor(0x2A, 0x2A, 0x2A), RGBColor(0x22, 0x22, 0x22)]
for r, row_data in enumerate(ret_data):
    bg = row_bg[r % 2]
    for c, val in enumerate(row_data):
        _style_cell(ret_table.cell(r + 1, c), val, bg=bg, fg=WHITE)

_add_text_box(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.4),
              "\u26a0 Auto-upgrades for gpt-4o Standard begin Mar 9, 2026. Review your deployments & test now.",
              font_size=13, color=RGBColor(0xFF, 0xA5, 0x00))


# ── Slide 9: Deprecation Action Items ───────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide, DARK_BG)
_add_rect(slide, 0, 0, W, Inches(0.15), ACCENT)
_add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
              "Retirement Action Items", font_size=34, bold=True, color=WHITE)
_accent_bar(slide, Inches(1.1))

items = [
    ("Immediate: Migrate gpt-5-chat Previews",
     "gpt-5-chat preview retires Mar 1, 2026. Switch to gpt-5.2-chat now."),
    ("By Mar 9: Test gpt-4o & gpt-4o-mini Workloads",
     "Auto-upgrades begin this date. Validate application behaviour on the latest GA models."),
    ("By Mar 31: Complete Standard Migrations",
     "All Standard gpt-4o / gpt-4o-mini deployments must be migrated or tested."),
    ("Migrate Fine-Tuned Models",
     "Fine-tuned deployments on retired bases are NOT auto-upgraded. "
     "Re-train on a supported base model before the retirement date."),
    ("Monitor the Retirements Page",
     "Bookmark learn.microsoft.com/azure/ai-foundry/openai/concepts/model-retirements."),
]
_add_bullet_slide(slide, items, Inches(1.5))


# ── Slide 10: February 2026 Quick-Reference Table ────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide, DARK_BG)
_add_rect(slide, 0, 0, W, Inches(0.15), ACCENT)
_add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
              "Key February 2026 Updates \u2013 Quick Reference", font_size=34, bold=True, color=WHITE)
_accent_bar(slide, Inches(1.1))

# Table
rows, cols = 12, 4
tbl_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5),
                                    Inches(12.3), Inches(5.4))
table = tbl_shape.table

# Column widths
table.columns[0].width = Inches(3.5)
table.columns[1].width = Inches(2.0)
table.columns[2].width = Inches(5.0)
table.columns[3].width = Inches(1.8)

headers = ["Item", "Category", "Highlight", "Date"]
data = [
    ["Claude Sonnet 4.6", "Model Release", "Frontier performance for coding & agents at scale", "Feb 17"],
    ["Claude Opus 4.6", "Model Release", "End-to-end autonomous task execution", "Feb 5"],
    ["PostgreSQL for AI", "Data Service", "DiskANN, MCP server, Copilot SQL, HorizonDB preview", "Feb 2"],
    ["High-Temp Superconductors", "Infrastructure", "Exploring superconductors for datacenter power", "Feb 10"],
    ["Agentic Cloud Ops", "Framework", "AI agents for cloud operations & incident response", "Feb 11"],
    ["Pantone Case Study", "Ecosystem", "Agentic AI + AI-ready database in production", "Feb 12"],
    ["Agents League", "Community", "Developer challenge for building AI agents", "Feb 16\u201327"],
    ["Agent Framework RC", "Foundry Blog", "Successor to SK & AutoGen hits Release Candidate", "Feb 19"],
    ["Foundry What\u2019s New", "Foundry Blog", "GPT-5.2, Codex Max, A2A, Memory, MCP roundup", "Feb 18"],
    ["DPO Fine-Tuning Guide", "Foundry Blog", "DPO walkthrough using Foundry SDK", "Feb 13"],
    ["Fine-Tuning Guide", "Foundry Blog", "When & how to move beyond prompting", "Feb 6"],
]

for i, h in enumerate(headers):
    _style_cell(table.cell(0, i), h, bold=True, bg=ACCENT, fg=WHITE, size=14)

row_bg = [RGBColor(0x2A, 0x2A, 0x2A), RGBColor(0x22, 0x22, 0x22)]
for r, row_data in enumerate(data):
    bg = row_bg[r % 2]
    for c, val in enumerate(row_data):
        _style_cell(table.cell(r + 1, c), val, bg=bg, fg=WHITE)


# ── Slide 11: Key Links ─────────────────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide, DARK_BG)
_add_rect(slide, 0, 0, W, Inches(0.15), ACCENT)
_add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
              "Key Links & Resources", font_size=34, bold=True, color=WHITE)
_accent_bar(slide, Inches(1.1))

links = [
    ("Microsoft Foundry Portal", "https://ai.azure.com/"),
    ("Foundry Product Page", "https://azure.microsoft.com/en-us/products/ai-foundry/"),
    ("Foundry Dev Blog", "https://devblogs.microsoft.com/foundry/"),
    ("Azure OpenAI What\u2019s New", "https://learn.microsoft.com/en-us/azure/ai-foundry/openai/whats-new"),
    ("Model Retirements", "https://learn.microsoft.com/en-us/azure/ai-foundry/openai/concepts/model-retirements"),
    ("Azure AI Blog", "https://azure.microsoft.com/en-us/blog/category/ai-machine-learning/"),
    ("Foundry Tech Community", "https://techcommunity.microsoft.com/category/azure-ai-foundry/blog/azure-ai-foundry-blog/"),
    ("Foundry Models Catalog", "https://azure.microsoft.com/en-us/products/ai-foundry/models/"),
    ("Foundry Agent Service", "https://azure.microsoft.com/en-us/products/ai-foundry/agent-service/"),
    ("Azure Updates Feed", "https://azure.microsoft.com/en-us/updates/"),
]

y = Inches(1.6)
for name, url in links:
    _add_text_box(slide, Inches(1.0), y, Inches(4), Inches(0.4), name,
                  font_size=16, bold=True, color=ACCENT)
    _add_text_box(slide, Inches(5.2), y, Inches(7.5), Inches(0.4), url,
                  font_size=14, color=LIGHT_GRAY)
    y += Inches(0.55)


# ── Slide 12: Closing ───────────────────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide, DARK_BG)
_add_rect(slide, 0, 0, W, Inches(0.15), ACCENT)

_add_text_box(slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.2),
              "Start Building with Microsoft Foundry",
              font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
_add_text_box(slide, Inches(0.8), Inches(3.8), Inches(11.5), Inches(0.8),
              "https://ai.azure.com",
              font_size=24, color=ACCENT, alignment=PP_ALIGN.CENTER)
_add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.6),
              "Compiled from official Microsoft sources \u2013 February 2026",
              font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────────────

output = r"c:\Users\petfue\Repos\whats-new-in-microsoft-ai\whats-new-microsoft-ai-feb2026.pptx"
prs.save(output)
print(f"Saved: {output}")
