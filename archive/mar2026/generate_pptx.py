"""Generate a PowerPoint deck summarizing Microsoft AI updates (Mar 1 - Apr 14, 2026)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DARK_BG = RGBColor(0x1B, 0x1B, 0x1B)
ACCENT = RGBColor(0x00, 0x78, 0xD4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width


def _set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=18,
    bold=False,
    color=WHITE,
    alignment=PP_ALIGN.LEFT,
    font_name="Segoe UI",
    hyperlink=None,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    p.alignment = alignment
    if hyperlink:
        run.hyperlink.address = hyperlink
    return box


def _accent_bar(slide, top):
    _add_rect(slide, Inches(0.5), top, Inches(1.2), Pt(4), ACCENT)


def _add_bullets(slide, items, top):
    box = slide.shapes.add_textbox(Inches(0.8), top, Inches(11.8), Inches(5.4))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        title = item[0]
        desc = item[1]
        url = item[2] if len(item) > 2 else None

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(8)
        p.space_after = Pt(4)

        run_title = p.add_run()
        run_title.text = title
        run_title.font.size = Pt(18)
        run_title.font.bold = True
        run_title.font.color.rgb = ACCENT
        run_title.font.name = "Segoe UI"
        if url:
            run_title.hyperlink.address = url

        run_desc = p.add_run()
        run_desc.text = "\n" + desc
        run_desc.font.size = Pt(15)
        run_desc.font.color.rgb = WHITE
        run_desc.font.name = "Segoe UI"


def _style_cell(cell, text, bold=False, bg=None, fg=WHITE, size=12, hyperlink=None):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = fg
    run.font.name = "Segoe UI"
    if hyperlink:
        run.hyperlink.address = hyperlink
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide, DARK_BG)
_add_rect(slide, 0, 0, W, Inches(0.15), ACCENT)
_add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1.2), "What's New in Microsoft AI", 44, True)
_add_text_box(
    slide,
    Inches(0.8),
    Inches(2.8),
    Inches(11.5),
    Inches(0.8),
    "Microsoft Foundry, Azure OpenAI Service & Related Platforms",
    22,
    False,
    LIGHT_GRAY,
)
_add_text_box(slide, Inches(0.8), Inches(4.0), Inches(11), Inches(0.6), "March 1 - April 14, 2026", 20, False, ACCENT)
_add_text_box(
    slide,
    Inches(0.8),
    Inches(6.2),
    Inches(11.5),
    Inches(0.5),
    "Sources: Microsoft Foundry Blog, Azure Blog, Microsoft Learn, Tech Community",
    13,
    False,
    LIGHT_GRAY,
)

# Slide 2: Major updates
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide, DARK_BG)
_add_rect(slide, 0, 0, W, Inches(0.15), ACCENT)
_add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8), "Major Platform Updates", 34, True)
_accent_bar(slide, Inches(1.1))
_add_bullets(
    slide,
    [
        (
            "What's New in Microsoft Foundry | March 2026 (Apr 9)",
            "Monthly roundup covering model releases, agent platform GA updates, safety and guardrails, platform operations, and SDK stabilization.",
            "https://devblogs.microsoft.com/foundry/whats-new-in-microsoft-foundry-mar-2026/",
        ),
        (
            "Priority Processing (GA) (Mar 23)",
            "Priority lane for latency-sensitive inference with pay-as-you-go flexibility and predictable performance for interactive AI workloads.",
            "https://techcommunity.microsoft.com/blog/azure-ai-foundry-blog/announcing-priority-processing-in-microsoft-foundry-for-performance-sensitive-ai/4504788",
        ),
    ],
    Inches(1.5),
)

# Slide 3: Models
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide, DARK_BG)
_add_rect(slide, 0, 0, W, Inches(0.15), ACCENT)
_add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8), "Model Releases & Catalog Expansion", 34, True)
_accent_bar(slide, Inches(1.1))
_add_bullets(
    slide,
    [
        (
            "GPT-5.4 + GPT-5.4 Pro (Mar 5)",
            "GA in Foundry with stronger instruction adherence, more dependable multi-step execution, and production-oriented reliability.",
            "https://techcommunity.microsoft.com/blog/azure-ai-foundry-blog/introducing-gpt-5-4-in-microsoft-foundry/4499785",
        ),
        (
            "GPT-5.4 mini + GPT-5.4 nano (Mar 17)",
            "Lower-latency variants for high-throughput tasks including classification, extraction, routing, and lightweight agent loops.",
            "https://techcommunity.microsoft.com/blog/azure-ai-foundry-blog/introducing-open-ai%E2%80%99s-gpt%E2%80%915-4-mini-in-microsoft-foundry/4500569",
        ),
        (
            "Grok 4.2 (GA) (Mar 30)",
            "xAI's Grok 4.2 moved to GA in the Foundry model catalog.",
            "https://devblogs.microsoft.com/foundry/whats-new-in-microsoft-foundry-mar-2026/",
        ),
    ],
    Inches(1.5),
)

# Slide 4: Agents and SDKs
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide, DARK_BG)
_add_rect(slide, 0, 0, W, Inches(0.15), ACCENT)
_add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8), "Agent Platform & SDKs", 34, True)
_accent_bar(slide, Inches(1.1))
_add_bullets(
    slide,
    [
        (
            "Foundry Agent Service is GA (Mar 16)",
            "Responses API-based runtime, private networking, expanded MCP authentication, and GA evaluations with continuous monitoring.",
            "https://devblogs.microsoft.com/foundry/foundry-agent-service-ga/",
        ),
        (
            "Hosted agents in additional regions",
            "Expansion includes East US, North Central US, Sweden Central, Southeast Asia, Japan East, and more.",
            "https://devblogs.microsoft.com/foundry/foundry-agent-service-ga/",
        ),
        (
            "SDK 2.0 stabilization across languages",
            "Python, JS/TS, and Java stable in March; .NET 2.0 GA on April 1. Unified AIProjectClient experience across workflows.",
            "https://devblogs.microsoft.com/foundry/whats-new-in-microsoft-foundry-mar-2026/",
        ),
    ],
    Inches(1.5),
)

# Slide 5: Open models and local
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide, DARK_BG)
_add_rect(slide, 0, 0, W, Inches(0.15), ACCENT)
_add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8), "Open Model Ecosystem & Edge", 34, True)
_accent_bar(slide, Inches(1.1))
_add_bullets(
    slide,
    [
        (
            "Fireworks AI on Foundry (Public Preview) (Mar 11)",
            "Open model inference integration with support for DeepSeek V3.2, gpt-oss-120b, Kimi K2.5, and MiniMax M2.5 plus BYOW.",
            "https://azure.microsoft.com/en-us/blog/introducing-fireworks-ai-on-microsoft-foundry-bringing-high-performance-low-latency-open-model-inference-to-azure/",
        ),
        (
            "NVIDIA Nemotron models added (Mar 16)",
            "Catalog expansion announced at GTC, broadening open model choices in Foundry.",
            "https://devblogs.microsoft.com/foundry/whats-new-in-microsoft-foundry-mar-2026/",
        ),
        (
            "Foundry Local is GA (Apr 9)",
            "Cross-platform local AI runtime for Windows, Linux, and macOS with offline inference and no per-token cloud cost.",
            "https://devblogs.microsoft.com/foundry/foundry-local-ga/",
        ),
    ],
    Inches(1.5),
)

# Slide 6: Lifecycle table
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide, DARK_BG)
_add_rect(slide, 0, 0, W, Inches(0.15), ACCENT)
_add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8), "Model Lifecycle Events (Mar 1 - Apr 14)", 34, True)
_accent_bar(slide, Inches(1.1))

rows, cols = 7, 4
tbl = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.6), Inches(12.3), Inches(4.9)).table

tbl.columns[0].width = Inches(3.8)
tbl.columns[1].width = Inches(2.2)
tbl.columns[2].width = Inches(2.0)
tbl.columns[3].width = Inches(4.3)

headers = ["Model / Family", "Event", "Date", "Notes"]
for i, h in enumerate(headers):
    _style_cell(tbl.cell(0, i), h, bold=True, bg=ACCENT, size=14)

data = [
    ["gpt-4o Standard", "Auto-upgrade start", "Mar 9", "Upgrade target: gpt-5.1"],
    ["gpt-4o-mini Standard", "Auto-upgrade start", "Mar 9", "Upgrade target: gpt-4.1-mini"],
    ["gpt-4o Standard", "Retired (Standard)", "Mar 31", "Other deployment types moved to Oct 1"],
    ["gpt-4o-mini Standard", "Retired (Standard)", "Mar 31", "Other deployment types moved to Oct 1"],
    ["gpt-4.1 / mini / nano", "Deprecated", "Apr 14", "Retirement date: Oct 14, 2026"],
    ["gpt-5-chat previews", "Retirement", "Apr 15", "Upgrade path: gpt-5.3-chat"],
]

row_bg = [RGBColor(0x2A, 0x2A, 0x2A), RGBColor(0x22, 0x22, 0x22)]
for r, row_data in enumerate(data):
    bg = row_bg[r % 2]
    for c, val in enumerate(row_data):
        _style_cell(tbl.cell(r + 1, c), val, bg=bg, size=12)

_add_text_box(
    slide,
    Inches(0.5),
    Inches(6.7),
    Inches(12.3),
    Inches(0.4),
    "Reference: learn.microsoft.com/en-us/azure/foundry/openai/concepts/model-retirements",
    12,
    False,
    ORANGE,
    hyperlink="https://learn.microsoft.com/en-us/azure/foundry/openai/concepts/model-retirements",
)

# Slide 7: Key actions
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide, DARK_BG)
_add_rect(slide, 0, 0, W, Inches(0.15), ACCENT)
_add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8), "Recommended Action Items", 34, True)
_accent_bar(slide, Inches(1.1))
_add_bullets(
    slide,
    [
        (
            "Confirm migration completion post-Mar 31",
            "Validate traffic and behavior after Standard gpt-4o / gpt-4o-mini retirement events.",
        ),
        (
            "Stop new gpt-4.1 family deployments",
            "As of Apr 14 deprecation, plan replacements on gpt-5 family alternatives.",
        ),
        (
            "Prepare for Apr 15 preview retirements",
            "Check any usage of gpt-5-chat previews and move to supported chat variants.",
        ),
        (
            "Plan o3 and o4-mini transitions",
            "Deprecation starts Apr 16, 2026. Build migration validation windows now.",
        ),
        (
            "Operationalize health alerts",
            "Use Azure Service Health advisories for upgrade, deprecation, and retirement notices.",
        ),
    ],
    Inches(1.5),
)

# Slide 8: Closing
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(slide, DARK_BG)
_add_rect(slide, 0, 0, W, Inches(0.15), ACCENT)
_add_text_box(
    slide,
    Inches(0.8),
    Inches(2.5),
    Inches(11.5),
    Inches(1.0),
    "Microsoft Foundry Momentum: Cloud to Edge",
    40,
    True,
    WHITE,
    PP_ALIGN.CENTER,
)
_add_text_box(
    slide,
    Inches(0.8),
    Inches(3.8),
    Inches(11.5),
    Inches(0.8),
    "https://ai.azure.com",
    24,
    False,
    ACCENT,
    PP_ALIGN.CENTER,
    hyperlink="https://ai.azure.com",
)
_add_text_box(
    slide,
    Inches(0.8),
    Inches(5.5),
    Inches(11.5),
    Inches(0.6),
    "Coverage period: Mar 1 - Apr 14, 2026",
    14,
    False,
    LIGHT_GRAY,
    PP_ALIGN.CENTER,
)

output = r"c:\Users\petfue\Repos\whats-new-in-microsoft-ai\whats-new-microsoft-ai-mar2026.pptx"
prs.save(output)
print(f"Saved: {output}")
